import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from fpdf import FPDF
from PIL import Image
import pytesseract

def convert_file(file_path: str):
    ext = os.path.splitext(file_path)[1].lower()
    output_file = None

    try:
        if ext == ".pdf":
            # PDF → TXT
            reader = PdfReader(file_path)
            text = "\n".join(page.extract_text() for page in reader.pages)
            output_file = file_path.replace(".pdf", ".txt")
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(text)

        elif ext == ".docx":
            # DOCX → PDF
            doc = Document(file_path)
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            for para in doc.paragraphs:
                pdf.multi_cell(0, 10, para.text)
            output_file = file_path.replace(".docx", ".pdf")
            pdf.output(output_file)

        elif ext == ".pptx":
            # PPTX → TXT
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            output_file = file_path.replace(".pptx", ".txt")
            with open(output_file, "w", encoding="utf-8") as f:
                f.write("\n".join(text))

        elif ext in [".png", ".jpg", ".jpeg"]:
            # Image → Text
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            output_file = file_path + ".txt"
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(text)

        return output_file
    except Exception as e:
        print("❌ Conversion error:", e)
        return None
